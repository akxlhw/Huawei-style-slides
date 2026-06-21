"""S4a: Render content.pptx.json to deck.pptx via MckEngine with Huawei theme."""

import sys
from pathlib import Path
from typing import Any, Dict, List, Tuple

MCK_ROOT = Path(__file__).resolve().parents[2] / "Mck-ppt-design-skill"
if str(MCK_ROOT) not in sys.path:
    sys.path.insert(0, str(MCK_ROOT))

from pptx.dml.color import RGBColor  # noqa: E402


# ──────────────────────────────────────────────────────────────────────────────
# Huawei visual spec applied to MckEngine constants before MckEngine is used.
# MckEngine hard-codes McKinsey navy/white; we monkey-patch its constants
# module so the generated deck follows the Huawei red/black/gray palette.
# ──────────────────────────────────────────────────────────────────────────────


def _hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert '#CF0A2C' to python-pptx RGBColor."""
    hex_color = hex_color.lstrip("#")
    return RGBColor(
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16),
    )


def _apply_huawei_theme() -> None:
    """Override MckEngine colors/fonts with the Huawei visual spec.

    MckEngine uses ``from .constants import *`` inside engine.py and core.py.
    That import copies the constant names into those modules at load time, so
    patching ``mck_ppt.constants`` alone is not enough: we also patch the
    engine and core module globals, plus the default argument values of any
    helper function that captured the original McKinsey RGBColor objects
    (e.g. ``add_oval(bg=NAVY)``).
    """
    import mck_ppt.constants as mc
    import mck_ppt.engine as me
    import mck_ppt.core as mcore

    theme = {
        # Primary palette — Huawei red is #CF0A2C (Pantone 185C, cross-verified)
        "NAVY": _hex_to_rgb("#CF0A2C"),       # Huawei red as the brand accent
        "BLACK": _hex_to_rgb("#231815"),      # Title black
        "WHITE": _hex_to_rgb("#FFFFFF"),
        # Gray scale
        "DARK_GRAY": _hex_to_rgb("#595757"),  # Body text
        "MED_GRAY": _hex_to_rgb("#999999"),   # Secondary text
        "LINE_GRAY": _hex_to_rgb("#DCDDDD"),  # Borders/dividers
        "BG_GRAY": _hex_to_rgb("#F5F5F5"),    # Table zebra stripe / light bg
        # Accent palette (Huawei report palette)
        "ACCENT_BLUE": _hex_to_rgb("#007DFF"),   # tech blue (Huawei Cloud primary)
        "ACCENT_GREEN": _hex_to_rgb("#669900"),  # warm green
        "ACCENT_ORANGE": _hex_to_rgb("#FF9900"), # warm orange
        "ACCENT_RED": _hex_to_rgb("#CF0A2C"),    # Huawei red
        "ACCENT_YELLOW": _hex_to_rgb("#FCC800"),
        # Light accent backgrounds
        "LIGHT_BLUE": _hex_to_rgb("#E6F0FF"),
        "LIGHT_GREEN": _hex_to_rgb("#EEF6E0"),
        "LIGHT_ORANGE": _hex_to_rgb("#FFF3E0"),
        "LIGHT_RED": _hex_to_rgb("#FCE8EA"),
        # Fonts
        "FONT_HEADER": "Microsoft YaHei",
        "FONT_BODY": "Microsoft YaHei",
        "FONT_EA": "Microsoft YaHei",
    }

    # Capture original colors BEFORE overwriting so we can patch function defaults.
    _old_to_new = {}
    for name, new_value in theme.items():
        if isinstance(new_value, RGBColor) and hasattr(mc, name):
            _old_to_new[getattr(mc, name)] = new_value

    for name, value in theme.items():
        setattr(mc, name, value)
        if hasattr(me, name):
            setattr(me, name, value)
        if hasattr(mcore, name):
            setattr(mcore, name, value)

    if not _old_to_new:
        return

    def _patch_module_defaults(mod):
        import inspect
        for _, obj in inspect.getmembers(mod, inspect.isfunction):
            if getattr(obj, "__defaults__", None):
                new_defaults = tuple(
                    _old_to_new.get(d, d) for d in obj.__defaults__
                )
                if new_defaults != obj.__defaults__:
                    obj.__defaults__ = new_defaults
            if getattr(obj, "__kwdefaults__", None):
                new_kwdefaults = {
                    k: _old_to_new.get(v, v) for k, v in obj.__kwdefaults__.items()
                }
                if new_kwdefaults != obj.__kwdefaults__:
                    obj.__kwdefaults__ = new_kwdefaults

    _patch_module_defaults(me)
    _patch_module_defaults(mcore)

    # ── Image placeholder → real image upgrade ─────────────────────────────
    # MckEngine's 7 image layouts (content_right_image, three_images, etc.)
    # all call add_image_placeholder(label='...') which draws a gray box.
    # We wrap it so that when `label` is a real file path (and the file
    # exists), a true picture is inserted via add_picture instead. This lets
    # content slides carry real images without changing MckEngine source.
    import os as _os
    _orig_img_ph = mcore.add_image_placeholder

    def _huawei_image(slide, left, top, width, height, label='Image'):
        if isinstance(label, str) and _os.path.isfile(label):
            slide.shapes.add_picture(label, left, top, width, height)
        else:
            _orig_img_ph(slide, left, top, width, height, label)

    mcore.add_image_placeholder = _huawei_image
    me.add_image_placeholder = _huawei_image

    # ── Arc-red title underline: OFF by default ─────────────────────────────
    # An arc/short-red-rule under the title was previously forced on every
    # page, but it read as odd/repetitive across a full deck. It is now
    # opt-in: a content slide can request it via slide["show_arc"] = True,
    # which _render_arc_underline() honors. The global add_action_title is left
    # unchanged (no per-page injection).


# Apply theme *before* importing MckEngine so its methods pick up Huawei colors.
_apply_huawei_theme()
from mck_ppt import MckEngine  # noqa: E402


def _convert_matrix_quadrants(
    quadrants: List[Tuple[str, str, str]],
) -> List[Tuple[str, RGBColor, str]]:
    """Convert (label, hex_color, desc) -> (label, RGBColor, desc)."""
    return [
        (str(q[0]), _hex_to_rgb(q[1]), str(q[2]))
        for q in quadrants
    ]


def _convert_axis_labels(axis_labels: Dict[str, str]) -> Tuple[str, str]:
    """Convert {'x': ..., 'y': ...} -> (x_label, y_label)."""
    if not axis_labels:
        return ("", "")
    return (axis_labels.get("x", ""), axis_labels.get("y", ""))


def _render_huawei_ring(eng, slide: Dict[str, Any]) -> None:
    """Render a native concentric-ring ecosystem diagram on the PPTX track.

    Draws multiple nested rings (outer→inner) using python-pptx OVAL shapes,
    a Huawei-red center disk, node dots + labels on each ring, and side cards.
    This replaces the earlier value_chain fallback so the PPTX track matches
    the HTML track's concentric-ring visual.
    """
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    import mck_ppt.core as mcore
    import mck_ppt.constants as mc

    s = eng._ns()
    mcore.add_action_title(s, slide["title"])

    rings = slide.get("rings", [])
    n = len(rings)
    # Ring geometry: center of canvas, outer radius fills ~38% of slide height.
    cx = Inches(4.3)            # ring center x (left half of content area)
    cy = Inches(4.25)           # ring center y
    outer_r = Inches(2.7)
    ring_gap = outer_r / max(n, 1)   # thickness per ring band

    # Draw rings outer→inner as stroked circles (line only, no fill).
    for ri, ring in enumerate(rings):
        r = outer_r - ring_gap * ri
        d = r * 2
        color = _hex_to_rgb(ring.get("color", "#CF0A2C"))
        oval = s.shapes.add_shape(MSO_SHAPE.OVAL, cx - r, cy - r, d, d)
        oval.fill.background()
        oval.line.color.rgb = color
        oval.line.width = Pt(1.5)
        mcore._clean_shape(oval)
        # (ring label is drawn in the consolidated loop below to avoid dupes)

    # Center disk (Huawei red, NAVY-patched).
    cd_r = Inches(0.55)
    center = s.shapes.add_shape(MSO_SHAPE.OVAL, cx - cd_r, cy - cd_r,
                                cd_r * 2, cd_r * 2)
    center.fill.solid()
    center.fill.fore_color.rgb = mc.NAVY
    center.line.fill.background()
    mcore._clean_shape(center)
    ctr = slide.get("center", {})
    mcore.add_text(s, cx - Inches(0.9), cy - Inches(0.2), Inches(1.8), Inches(0.4),
                   ctr.get("label", ""), font_size=Pt(14), font_color=mc.WHITE,
                   bold=True, alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Ring labels carry their items inline (no scattered node dots), which
    # avoids the overlap warnings that scattered nodes triggered across rings.
    import math
    for ri, ring in enumerate(rings):
        r = outer_r - ring_gap * ri
        color = _hex_to_rgb(ring.get("color", "#CF0A2C"))
        items = ring.get("items", [])
        # Combined label at the top of the ring band.
        combined = ring.get("label", "") + ("：" + " · ".join(items) if items else "")
        mcore.add_text(s, cx - Inches(1.5), cy - r - Inches(0.05),
                       Inches(3.0), Inches(0.3), combined,
                       font_size=Pt(10), font_color=color, bold=True,
                       alignment=PP_ALIGN.CENTER)

    # Side cards on the right half.
    sc_x = Inches(8.6)
    sc_y = Inches(1.6)
    sc_w = Inches(3.8)
    sc_h = Inches(1.4)
    for i, card in enumerate(slide.get("side_cards", [])):
        cy_i = sc_y + (sc_h + Inches(0.2)) * i
        title, color_hex, desc = card[0], card[1], card[2]
        color = _hex_to_rgb(color_hex)
        mcore.add_rect(s, sc_x, cy_i, sc_w, sc_h, mc.BG_GRAY)
        mcore.add_rect(s, sc_x, cy_i, Inches(0.08), sc_h, color)
        mcore.add_text(s, sc_x + Inches(0.25), cy_i + Inches(0.15),
                       sc_w - Inches(0.4), Inches(0.35), str(title),
                       font_size=Pt(15), font_color=color, bold=True)
        mcore.add_text(s, sc_x + Inches(0.25), cy_i + Inches(0.55),
                       sc_w - Inches(0.4), sc_h - Inches(0.6), str(desc),
                       font_size=Pt(11), font_color=mc.DARK_GRAY)

    eng._footer(s, slide.get("source", ""))


def _render_huawei_nav(eng, slide: Dict[str, Any], chapters: list, current_chapter: str) -> None:
    """Post-process: draw a slim left nav rail on an already-rendered slide.

    Drawn AFTER the slide body so it sits on top. Uses a narrow 0.55\" band at
    the far left edge so it does not collide with content (which starts at
    LM=0.8\"). Shows the current chapter in Huawei red; others in muted gray.
    """
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import MSO_ANCHOR
    import mck_ppt.core as mcore
    import mck_ppt.constants as mc

    if not chapters:
        return
    # The last slide rendered is eng.prs.slides[-1] (eng just rendered it).
    s = eng.prs.slides[-1]
    nav_w = Inches(0.55)
    # Background band
    mcore.add_rect(s, Emu(0), Inches(0.0), nav_w, Inches(7.5), mc.BG_GRAY)
    # Brand tick
    mcore.add_rect(s, Emu(0), Inches(0.0), Inches(0.06), Inches(7.5), mc.NAVY)
    # Chapters stacked vertically, small text. Each label gets a FIXED small
    # height (not slot_h) so the text boxes don't span the full vertical area
    # and trigger shape_overlap warnings against content-area text.
    n = len(chapters)
    if n == 0:
        return
    slot_h = Inches(6.0) / n
    start_y = Inches(1.2)
    label_h = Inches(0.28)   # fixed short height per chapter label
    label_w = Inches(0.34)
    for i, ch in enumerate(chapters):
        y = start_y + slot_h * i
        is_active = (ch == current_chapter)
        color = mc.NAVY if is_active else mc.MED_GRAY
        bold = is_active
        # Active marker dot
        if is_active:
            from pptx.enum.shapes import MSO_SHAPE
            dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.12), y + Inches(0.04),
                                     Inches(0.10), Inches(0.10))
            dot.fill.solid(); dot.fill.fore_color.rgb = mc.NAVY
            dot.line.fill.background(); mcore._clean_shape(dot)
        # Chapter label — short fixed box, confined to the nav band (x<0.55").
        mcore.add_text(s, Inches(0.20), y, label_w, label_h,
                       str(ch), font_size=Pt(8), font_color=color, bold=bold,
                       anchor=MSO_ANCHOR.TOP)


# ── Native Huawei chart renderers (iteration round 2) ───────────────────────
# These draw directly with python-pptx shapes + MckEngine helpers because the
# MckEngine has no built-in swim-lane / true-pyramid / radar / gantt methods.
# All follow the standard pattern: eng._ns() → add_action_title → draw → _footer.


# Native renderers (swim-lane, radar, gantt, etc.) — extracted for clarity.
# _render_huawei_ring and _render_huawei_nav stay in this file (they use
# _convert_* helpers defined here); the rest are in native_renderers.py.
from native_renderers import (
    _new_slide,
    _render_swim_lane, _render_pyramid_triangle, _render_radar_chart,
    _render_gantt_chart, _render_icon_stat_grid, _render_cycle_ring,
    _render_image_bg,
)


def render_pptx(content: Dict[str, Any], output_dir: Path) -> Path:
    """Render a content.pptx.json dict to deck.pptx using MckEngine.

    Includes a simple content-hash cache: if the deck already exists and the
    content hasn't changed (same hash), skip re-rendering to save time.
    """
    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)
    deck_path = output_dir / "deck.pptx"

    # ── Cache check: skip render if content unchanged ──────────────────────
    import hashlib
    import json as _json
    content_hash = hashlib.md5(
        _json.dumps(content, ensure_ascii=False, sort_keys=True).encode()
    ).hexdigest()[:12]
    hash_file = output_dir / ".deck_content_hash"
    if deck_path.is_file() and hash_file.is_file():
        if hash_file.read_text().strip() == content_hash:
            return deck_path  # cache hit — skip rendering

    total = len(content["slides"])
    eng = MckEngine(total_slides=total)

    # Chapter metadata for the Huawei left-nav rail (mirrors the HTML track).
    chapters = content.get("chapters", [])
    nav_scenarios = ("skill-report",)
    use_nav = content.get("scenario") in nav_scenarios and bool(chapters)

    for slide in content["slides"]:
        layout = slide["layout"]
        if layout == "cover":
            cover_image = slide.get("cover_image")
            eng.cover(
                title=slide["title"],
                subtitle=slide.get("subtitle", ""),
                author=slide.get("author", ""),
                date=slide.get("date", ""),
                cover_image=cover_image if cover_image else None,
            )
        elif layout == "executive_summary":
            eng.executive_summary(
                title=slide["title"],
                headline=slide.get("headline", ""),
                items=slide.get("items", []),
                source=slide.get("source", ""),
            )
        elif layout == "table_insight":
            eng.table_insight(
                title=slide["title"],
                headers=slide["headers"],
                rows=slide["rows"],
                insights=slide.get("insights", []),
                source=slide.get("source", ""),
            )
        elif layout == "funnel":
            # The retired eng.funnel() overflows; use process_chevron instead.
            steps = [
                (str(i + 1), str(label), f"{count} ({int(pct * 100)}%)")
                for i, (label, count, pct) in enumerate(slide.get("stages", []))
            ]
            eng.process_chevron(
                title=slide["title"],
                steps=steps,
                source=slide.get("source", ""),
            )
        elif layout == "data_table":
            eng.data_table(
                title=slide["title"],
                headers=slide["headers"],
                rows=slide["rows"],
                source=slide.get("source", ""),
            )
        elif layout == "matrix_2x2":
            axis_labels = _convert_axis_labels(slide.get("axis_labels"))
            eng.matrix_2x2(
                title=slide["title"],
                quadrants=_convert_matrix_quadrants(slide.get("quadrants", [])),
                axis_labels=axis_labels if any(axis_labels) else None,
                source=slide.get("source", ""),
            )
        elif layout == "timeline":
            eng.timeline(
                title=slide["title"],
                milestones=slide.get("milestones", []),
                source=slide.get("source", ""),
            )
        elif layout == "action_items":
            eng.action_items(
                title=slide["title"],
                actions=slide.get("actions", []),
                source=slide.get("source", ""),
            )
        elif layout == "closing":
            eng.closing(
                title=slide["title"],
                message=slide.get("message", ""),
            )
        elif layout == "big_number":
            eng.big_number(
                title=slide["title"],
                number=slide.get("number", ""),
                unit=slide.get("unit", ""),
                description=slide.get("description", ""),
                detail_items=slide.get("detail_items"),
                source=slide.get("source", ""),
            )
        elif layout == "three_stat":
            stats = [(s[0], s[1], bool(s[2])) for s in slide.get("stats", [])]
            eng.three_stat(
                title=slide["title"],
                stats=stats,
                detail_items=slide.get("detail_items"),
                source=slide.get("source", ""),
            )
        elif layout == "metric_cards":
            cards = []
            for card in slide.get("cards", []):
                if len(card) == 5:
                    letter, ctitle, desc, accent, light = card
                    cards.append((letter, ctitle, desc, _hex_to_rgb(accent), _hex_to_rgb(light)))
                else:
                    cards.append(tuple(card[:3]))
            eng.metric_cards(
                title=slide["title"],
                cards=cards,
                source=slide.get("source", ""),
            )
        elif layout == "grouped_bar":
            series = [(n, _hex_to_rgb(c)) for n, c in slide.get("series", [])]
            eng.grouped_bar(
                title=slide["title"],
                categories=slide.get("categories", []),
                series=series,
                data=slide.get("data", []),
                y_ticks=slide.get("y_ticks"),
                summary=slide.get("summary"),
                source=slide.get("source", ""),
            )
        elif layout == "horizontal_bar":
            items = [(n, p, _hex_to_rgb(c)) for n, p, c in slide.get("items", [])]
            eng.horizontal_bar(
                title=slide["title"],
                items=items,
                summary=slide.get("summary"),
                source=slide.get("source", ""),
            )
        elif layout == "donut_chart":
            segments = [(p, _hex_to_rgb(c), lab)
                        for p, c, lab in slide.get("segments", [])]
            eng.donut(
                title=slide["title"],
                segments=segments,
                center_label=slide.get("center_label", ""),
                center_sub=slide.get("center_sub", ""),
                summary=slide.get("summary"),
                source=slide.get("source", ""),
            )
        elif layout == "kpi_tracker":
            eng.kpi_tracker(
                title=slide["title"],
                kpis=slide.get("kpis", []),
                summary=slide.get("summary"),
                source=slide.get("source", ""),
            )
        elif layout == "pyramid_steps":
            levels = [(lab, desc, icon) for lab, desc, icon in slide.get("levels", [])]
            eng.pyramid(
                title=slide["title"],
                levels=levels,
                source=slide.get("source", ""),
            )
        elif layout == "four_column":
            eng.four_column(
                title=slide["title"],
                items=slide.get("items", []),
                source=slide.get("source", ""),
            )
        elif layout == "scorecard":
            eng.scorecard(
                title=slide["title"],
                items=slide.get("items", []),
                source=slide.get("source", ""),
            )
        elif layout == "checklist_status":
            from pptx.util import Inches as _Inches
            col_widths = [_Inches(float(w)) for w in slide.get("col_widths", [])]
            eng.checklist(
                title=slide["title"],
                columns=slide.get("columns", []),
                col_widths=col_widths,
                rows=[tuple(r) for r in slide.get("rows", [])],
                source=slide.get("source", ""),
            )
        elif layout == "value_chain":
            stages = [(lab, desc, _hex_to_rgb(c))
                      for lab, desc, c in slide.get("stages", [])]
            eng.value_chain(
                title=slide["title"],
                stages=stages,
                source=slide.get("source", ""),
            )
        elif layout == "ecosystem_ring":
            # Native concentric-ring diagram (Huawei signature shape) on PPTX.
            _render_huawei_ring(eng, slide)
        elif layout == "swot":
            quadrants = [
                (lab, _hex_to_rgb(accent), _hex_to_rgb(bg), pts)
                for lab, accent, bg, pts in slide.get("quadrants", [])
            ]
            eng.swot(
                title=slide["title"],
                quadrants=quadrants,
                source=slide.get("source", ""),
            )
        elif layout == "icon_grid":
            items = [(t, d, _hex_to_rgb(c)) for t, d, c in slide.get("items", [])]
            eng.icon_grid(
                title=slide["title"],
                items=items,
                source=slide.get("source", ""),
            )
        elif layout == "quote":
            # quote() has no title/source params
            eng.quote(
                quote_text=slide.get("quote_text", ""),
                attribution=slide.get("attribution", ""),
            )
        elif layout == "cycle":
            # Native ring layout (avoids MckEngine cycle's node/panel overlap).
            _render_cycle_ring(eng, slide)
        elif layout == "side_by_side":
            options = [(str(o[0]), list(o[1])) for o in slide.get("options", [])]
            eng.side_by_side(
                title=slide["title"], options=options, source=slide.get("source", ""),
            )
        elif layout == "case_study":
            sections = [(str(s[0]), str(s[1]), str(s[2]))
                        for s in slide.get("sections", [])]
            rb = slide.get("result_box")
            result_box = (str(rb[0]), str(rb[1])) if rb else None
            eng.case_study(
                title=slide["title"], sections=sections, result_box=result_box,
                source=slide.get("source", ""),
            )
        elif layout == "decision_tree":
            root = (str(slide["root"][0]),) if slide.get("root") else ("",)
            branches = [
                (str(b[0]), str(b[1]), _hex_to_rgb(b[2]),
                 [(str(c[0]), str(c[1])) for c in b[3]])
                for b in slide.get("branches", [])
            ]
            rp = slide.get("right_panel")
            right_panel = (str(rp[0]), list(rp[1])) if rp else None
            eng.decision_tree(
                title=slide["title"], root=root, branches=branches,
                right_panel=right_panel, source=slide.get("source", ""),
            )
        elif layout == "harvey_ball_compare":
            eng.harvey_ball_table(
                title=slide["title"],
                criteria=slide.get("criteria", []),
                options=slide.get("options", []),
                scores=slide.get("scores", []),
                legend_text=slide.get("legend_text"),
                summary=slide.get("summary"),
                source=slide.get("source", ""),
            )
        elif layout == "swim_lane":
            _render_swim_lane(eng, slide)
        elif layout == "pyramid_triangle":
            _render_pyramid_triangle(eng, slide)
        elif layout == "radar_chart":
            _render_radar_chart(eng, slide)
        elif layout == "gantt_chart":
            _render_gantt_chart(eng, slide)
        elif layout == "icon_stat_grid":
            _render_icon_stat_grid(eng, slide)
        else:
            raise ValueError(f"Unsupported PPTX layout: {layout}")

        # Full-bleed background image support (content slides with an `image`).
        if layout not in ("cover", "closing") and slide.get("image"):
            _render_image_bg(eng.prs.slides[-1], slide["image"])

        # Huawei left-nav rail (post-process, only on content slides for the
        # skill-report scenario so it matches the HTML track).
        if use_nav and layout not in ("cover", "closing"):
            _render_huawei_nav(eng, slide, chapters, slide.get("chapter", ""))

        # Optional arc-red title underline (off by default; a slide opts in
        # via slide["show_arc"] = True).
        if slide.get("show_arc") and layout not in ("cover", "closing"):
            from pptx.util import Inches as _In, Pt as _Pt
            import mck_ppt.core as _mc
            _mc.add_hline(eng.prs.slides[-1], mc.LM, mc.TITLE_LINE_Y,
                          _In(2.4), mc.NAVY, _Pt(2.5))

    deck_path = output_dir / "deck.pptx"
    try:
        eng.save(str(deck_path))
    except UnicodeEncodeError:
        # MckEngine's success print uses emoji that can fail on GBK/ASCII
        # consoles. The file is already saved at this point, so verify it
        # exists and continue.
        if not deck_path.exists():
            raise
    # Save content hash for cache on next run.
    hash_file.write_text(content_hash)
    return deck_path
