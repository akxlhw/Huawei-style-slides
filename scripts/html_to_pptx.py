"""HTML→PPTX converter: screenshot each HTML slide via Playwright and embed
the PNGs into a PPTX. This lets the Huawei-styled HTML deck become an
editable PPTX fallback (images on slides), complementing the native MckEngine
PPTX track.

Usage:
    from html_to_pptx import html_to_pptx
    html_to_pptx('slides.html', 'deck_from_html.pptx')

    # Or via CLI:
    python scripts/html_to_pptx.py slides.html -o deck_from_html.pptx
"""
import argparse
import asyncio
from pathlib import Path
from typing import Optional


def html_to_pptx(
    html_path: str,
    output_path: str,
    width: int = 1920,
    height: int = 1080,
) -> Path:
    """Convert a self-contained HTML deck to a PPTX by screenshotting slides.

    Each <section class="slide"> is rendered as a full-slide image in the PPTX.
    Requires Playwright (chromium) installed.
    """
    from pptx import Presentation
    from pptx.util import Inches

    html_path = Path(html_path)
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    # Step 1: screenshot each slide to PNG via Playwright.
    pngs = _screenshot_slides(html_path, output_path.parent, width, height)
    if not pngs:
        raise RuntimeError("No slides found in HTML — expected <section class=\"slide\">")

    # Step 2: build PPTX with one slide per PNG.
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    for png in pngs:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(png), 0, 0, prs.slide_width, prs.slide_height)

    prs.save(str(output_path))
    print(f"✅ HTML→PPTX: {len(pngs)} slides → {output_path}")
    return output_path


def _screenshot_slides(html_path: Path, out_dir: Path, width: int, height: int) -> list:
    """Screenshot each slide section to PNG. Returns list of Path."""
    try:
        from playwright.async_api import async_playwright
    except ImportError:
        raise RuntimeError(
            "Playwright is required for HTML→PPTX. Install: "
            "pip install playwright && playwright install chromium"
        )

    async def shoot():
        pngs = []
        async with async_playwright() as p:
            browser = await p.chromium.launch()
            page = await browser.new_page(viewport={"width": width, "height": height})
            await page.goto(f"file:///{html_path.resolve().as_posix()}")
            await page.wait_for_timeout(2000)

            # Find slide sections.
            slides = await page.query_selector_all('section.slide')
            if not slides:
                # Try deck-stage children.
                slides = await page.query_selector_all('[data-deck-active], deck-stage > *')

            for i in range(len(slides)):
                # Navigate to slide i (keyboard nav: Home then i×ArrowRight).
                await page.keyboard.press("Home")
                await page.wait_for_timeout(200)
                for _ in range(i):
                    await page.keyboard.press("ArrowRight")
                    await page.wait_for_timeout(100)
                await page.wait_for_timeout(400)
                png = out_dir / f"_htmlslide_{i:02d}.png"
                await page.screenshot(path=str(png))
                pngs.append(png)

            await browser.close()
        return pngs

    return asyncio.run(shoot())


def main():
    parser = argparse.ArgumentParser(description="Convert Huawei HTML deck to PPTX")
    parser.add_argument("html", help="path to slides.html")
    parser.add_argument("-o", "--output", default="deck_from_html.pptx",
                        help="output PPTX path")
    args = parser.parse_args()
    html_to_pptx(args.html, args.output)


if __name__ == "__main__":
    main()
