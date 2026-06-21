"""Native Huawei chart renderers — drawn directly with python-pptx.

Extracted from pptx_renderer.py for maintainability. These functions draw
shapes (swim-lane, radar, gantt, etc.) that MckEngine has no built-in for.
All follow: eng._ns() → add_action_title → draw → _footer.
"""
import math
from typing import Any, Dict, Tuple
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor


def _hex_to_rgb(hex_color):
    hex_color = hex_color.lstrip('#')
    return RGBColor(int(hex_color[0:2],16), int(hex_color[2:4],16), int(hex_color[4:6],16))


def _new_slide(eng, title, source):
    """Shared scaffold: create slide, draw Huawei action title, return (s, mcore, mc)."""
    import mck_ppt.core as mcore
    import mck_ppt.constants as mc
    s = eng._ns()
    mcore.add_action_title(s, title)
    return s, mcore, mc


def _render_swim_lane(eng, slide):
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    s, mcore, mc = _new_slide(eng, slide["title"], slide.get("source", ""))
    lanes = slide.get("lanes", [])
    phases = slide.get("phases", [])
    steps = slide.get("steps", [])
    n_lanes = len(lanes)
    n_phases = len(phases)
    if n_lanes == 0 or n_phases == 0:
        eng._footer(s, slide.get("source", "")); return
    grid_x = Inches(1.6)
    grid_y = Inches(1.6)
    lane_w = Inches(9.6)
    col_w = lane_w / n_phases
    lane_h = Inches(4.8) / n_lanes
    label_w = Inches(1.4)
    # Phase headers
    for j, ph in enumerate(phases):
        mcore.add_text(s, grid_x + col_w * j, grid_y - Inches(0.4), col_w, Inches(0.35),
                       str(ph), font_size=Pt(12), font_color=mc.NAVY, bold=True,
                       alignment=PP_ALIGN.CENTER)
    # Lane rows
    lane_colors = [mc.LIGHT_RED, mc.LIGHT_BLUE, mc.LIGHT_GREEN, mc.LIGHT_ORANGE]
    for i, lane in enumerate(lanes):
        ly = grid_y + lane_h * i
        bg = lane_colors[i % len(lane_colors)]
        mcore.add_rect(s, grid_x, ly, lane_w, lane_h, bg)
        mcore.add_text(s, Inches(0.2), ly, label_w, lane_h, str(lane),
                       font_size=Pt(12), font_color=mc.DARK_GRAY, bold=True,
                       anchor=MSO_ANCHOR.MIDDLE)
        # Phase column dividers
        if i == 0:
            for j in range(1, n_phases):
                mcore.add_rect(s, grid_x + col_w * j, ly, Pt(1), Inches(4.8), mc.LINE_GRAY)
    # Steps
    for (li, pi, label) in steps:
        if li >= n_lanes or pi >= n_phases:
            continue
        bx = grid_x + col_w * pi + Inches(0.15)
        by = grid_y + lane_h * li + Inches(0.12)
        bw = col_w - Inches(0.3)
        bh = lane_h - Inches(0.24)
        mcore.add_rect(s, bx, by, bw, bh, mc.WHITE)
        mcore.add_rect(s, bx, by, bw, Inches(0.05), mc.NAVY)
        mcore.add_text(s, bx + Inches(0.08), by + Inches(0.1), bw - Inches(0.16), bh - Inches(0.15),
                       str(label), font_size=Pt(10), font_color=mc.DARK_GRAY,
                       alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    eng._footer(s, slide.get("source", ""))


def _render_pyramid_triangle(eng, slide):
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    s, mcore, mc = _new_slide(eng, slide["title"], slide.get("source", ""))
    levels = slide.get("levels", [])
    n = len(levels)
    if n == 0:
        eng._footer(s, slide.get("source", "")); return
    cx = Inches(6.66)  # slide center-ish
    top_y = Inches(1.6)
    max_w = Inches(8.0)
    layer_h = Inches(0.95)
    for i, (label, desc, color_hex) in enumerate(levels):
        color = _hex_to_rgb(color_hex)
        # Trapezoid approximated by rectangle centered, decreasing width.
        w = max_w * (n - i) / n
        y = top_y + layer_h * i
        shape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx - w / 2, y, w, layer_h - Inches(0.08))
        shape.fill.solid(); shape.fill.fore_color.rgb = color
        shape.line.fill.background(); mcore._clean_shape(shape)
        mcore.add_text(s, cx - w / 2, y, w, layer_h - Inches(0.08),
                       f"{label} · {desc}", font_size=Pt(12), font_color=mc.WHITE,
                       bold=True, alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    eng._footer(s, slide.get("source", ""))


def _render_radar_chart(eng, slide):
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    s, mcore, mc = _new_slide(eng, slide["title"], slide.get("source", ""))
    axes = slide.get("axes", [])
    series = slide.get("series", [])
    n = len(axes)
    if n < 3:
        eng._footer(s, slide.get("source", "")); return
    import math
    cx = Inches(5.5); cy = Inches(4.55); R = Inches(2.25)
    # Concentric grid rings (3 levels)
    for frac in (0.33, 0.66, 1.0):
        d = R * 2 * frac
        ring = s.shapes.add_shape(MSO_SHAPE.OVAL, cx - R * frac, cy - R * frac, d, d)
        ring.fill.background(); ring.line.color.rgb = mc.LINE_GRAY; ring.line.width = Pt(0.5)
        mcore._clean_shape(ring)
    # Spokes + axis labels
    for i, ax in enumerate(axes):
        ang = math.radians(-90 + i * (360 / n))
        ex = cx + int(R * math.cos(ang)); ey = cy + int(R * math.sin(ang))
        mcore.add_rect(s, min(cx, ex), min(cy, ey), abs(ex - cx) or Pt(1), abs(ey - cy) or Pt(1), mc.LINE_GRAY)
        lx = cx + int(R * 1.18 * math.cos(ang)) - Inches(0.6)
        ly = cy + int(R * 1.18 * math.sin(ang)) - Inches(0.12)
        mcore.add_text(s, lx, ly, Inches(1.2), Inches(0.24), str(ax),
                       font_size=Pt(10), font_color=mc.DARK_GRAY, bold=True,
                       alignment=PP_ALIGN.CENTER)
    # Series data polygons
    for ser in series:
        color = _hex_to_rgb(ser.get("color", "#CF0A2C"))
        vals = ser.get("values", [])
        pts = []
        for i, v in enumerate(vals[:n]):
            ang = math.radians(-90 + i * (360 / n))
            r = R * (v / 100.0)
            pts.append((cx + int(r * math.cos(ang)), cy + int(r * math.sin(ang))))
        # Draw filled polygon via freeform
        if len(pts) >= 3:
            try:
                ff = s.shapes.build_freeform(pts[0][0], pts[0][1])
                ff.add_line_segments([(x, y) for x, y in pts[1:]], close=True)
                shape = ff.convert_to_shape()
                shape.fill.solid(); shape.fill.fore_color.rgb = color
                from pptx.dml.color import RGBColor
                shape.fill.fore_color.rgb = color
                # set transparency via alpha not directly supported; use line only overlay
                shape.line.color.rgb = color; shape.line.width = Pt(1.5)
                mcore._clean_shape(shape)
            except Exception:
                # Fallback: dots at each vertex
                for px, py in pts:
                    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, px - Inches(0.05), py - Inches(0.05),
                                             Inches(0.1), Inches(0.1))
                    dot.fill.solid(); dot.fill.fore_color.rgb = color
                    dot.line.fill.background(); mcore._clean_shape(dot)
    # Legend
    lgx = Inches(10.5)
    for i, ser in enumerate(series):
        color = _hex_to_rgb(ser.get("color", "#CF0A2C"))
        ly = Inches(2.4) + Inches(0.4) * i
        mcore.add_rect(s, lgx, ly, Inches(0.2), Inches(0.2), color)
        mcore.add_text(s, lgx + Inches(0.3), ly - Inches(0.02), Inches(2.0), Inches(0.28),
                       str(ser.get("name", "")), font_size=Pt(11), font_color=mc.DARK_GRAY)
    eng._footer(s, slide.get("source", ""))


def _render_gantt_chart(eng, slide):
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    s, mcore, mc = _new_slide(eng, slide["title"], slide.get("source", ""))
    phases = slide.get("phases", [])
    tasks = slide.get("tasks", [])
    n_ph = len(phases)
    n_tk = len(tasks)
    if n_ph == 0 or n_tk == 0:
        eng._footer(s, slide.get("source", "")); return
    gx = Inches(2.6); gy = Inches(1.7)
    track_w = Inches(8.4); col_w = track_w / n_ph
    row_h = Inches(4.6) / n_tk
    label_w = Inches(1.9)
    label_x = Inches(0.7)   # right of the nav rail (0-0.55"), avoid overlap
    # Phase headers
    for j, ph in enumerate(phases):
        mcore.add_text(s, gx + col_w * j, gy - Inches(0.4), col_w, Inches(0.32),
                       str(ph), font_size=Pt(11), font_color=mc.NAVY, bold=True,
                       alignment=PP_ALIGN.CENTER)
    # Vertical phase dividers
    for j in range(n_ph + 1):
        mcore.add_rect(s, gx + col_w * j, gy, Pt(1), Inches(4.6), mc.LINE_GRAY)
    # Tasks
    for i, (name, start, end, color_hex) in enumerate(tasks):
        ty = gy + row_h * i + Inches(0.08)
        mcore.add_text(s, label_x, ty, label_w, row_h - Inches(0.16), str(name),
                       font_size=Pt(11), font_color=mc.DARK_GRAY, anchor=MSO_ANCHOR.MIDDLE)
        bar_x = gx + col_w * start
        bar_w = col_w * (end - start)
        color = _hex_to_rgb(color_hex)
        bar = s.shapes.add_shape(__import__('pptx.enum.shapes', fromlist=['MSO_SHAPE']).MSO_SHAPE.ROUNDED_RECTANGLE,
                                 bar_x, ty + Inches(0.05), bar_w, row_h - Inches(0.2))
        bar.fill.solid(); bar.fill.fore_color.rgb = color
        bar.line.fill.background(); mcore._clean_shape(bar)
    eng._footer(s, slide.get("source", ""))


def _render_icon_stat_grid(eng, slide):
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    s, mcore, mc = _new_slide(eng, slide["title"], slide.get("source", ""))
    items = slide.get("items", [])
    n = len(items)
    if n == 0:
        eng._footer(s, slide.get("source", "")); return
    cols = 3
    rows = (n + cols - 1) // cols
    import math
    cell_w = Inches(9.6) / cols
    cell_h = Inches(4.4) / rows
    gx = Inches(1.6); gy = Inches(1.7)
    for i, (num, label, desc, color_hex) in enumerate(items):
        col = i % cols; row = i // cols
        x = gx + cell_w * col; y = gy + cell_h * row
        color = _hex_to_rgb(color_hex)
        mcore.add_rect(s, x + Inches(0.1), y + Inches(0.1), cell_w - Inches(0.2), cell_h - Inches(0.2), mc.BG_GRAY)
        mcore.add_rect(s, x + Inches(0.1), y + Inches(0.1), cell_w - Inches(0.2), Inches(0.06), color)
        # Big number
        mcore.add_text(s, x + Inches(0.1), y + Inches(0.25), cell_w - Inches(0.2), Inches(0.7),
                       str(num), font_size=Pt(36), font_color=color, bold=True,
                       alignment=PP_ALIGN.CENTER, font_name=mc.FONT_HEADER)
        mcore.add_text(s, x + Inches(0.1), y + Inches(0.95), cell_w - Inches(0.2), Inches(0.35),
                       str(label), font_size=Pt(13), font_color=mc.DARK_GRAY, bold=True,
                       alignment=PP_ALIGN.CENTER)
        mcore.add_text(s, x + Inches(0.2), y + Inches(1.3), cell_w - Inches(0.4), Inches(0.5),
                       str(desc), font_size=Pt(10), font_color=mc.DARK_GRAY,
                       alignment=PP_ALIGN.CENTER)
    eng._footer(s, slide.get("source", ""))


def _render_cycle_ring(eng, slide):
    """Native cycle: nodes in a ring (left half) + side panel (right half).

    Replaces MckEngine's cycle() whose horizontal node line collides with the
    right panel text box. Here nodes orbit a center hub on the left, and the
    panel sits in a clearly separated right column.
    """
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    import math
    s, mcore, mc = _new_slide(eng, slide["title"], slide.get("source", ""))

    phases = slide.get("phases", [])
    n = len(phases)
    # Ring on the left half: center at (4.2, 4.2), radius 2.2.
    cx, cy, R = Inches(4.2), Inches(4.25), Inches(2.2)
    # Center hub
    hub_r = Inches(0.6)
    hub = s.shapes.add_shape(MSO_SHAPE.OVAL, cx - hub_r, cy - hub_r, hub_r * 2, hub_r * 2)
    hub.fill.solid(); hub.fill.fore_color.rgb = mc.NAVY
    hub.line.fill.background(); mcore._clean_shape(hub)

    # Node boxes orbiting the hub.
    for i, ph in enumerate(phases):
        label = str(ph[0])
        ang = math.radians(-90 + i * (360 / max(n, 1)))
        nx = cx + int(R * math.cos(ang))
        ny = cy + int(R * math.sin(ang))
        bw, bh = Inches(1.5), Inches(0.7)
        bx, by = nx - bw // 2, ny - bh // 2
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, by, bw, bh)
        box.fill.solid(); box.fill.fore_color.rgb = mc.WHITE
        box.line.color.rgb = mc.NAVY; box.line.width = Pt(1.5)
        mcore._clean_shape(box)
        mcore.add_text(s, bx, by, bw, bh, label, font_size=Pt(14),
                       font_color=mc.NAVY, bold=True, alignment=PP_ALIGN.CENTER,
                       anchor=MSO_ANCHOR.MIDDLE)
        # Arrow hint to next node (curved arrow as text glyph near the edge)
        # omitted to keep geometry clean — the ring order is implied by layout.

    # Side panel on the right half (x ≥ 7.5"), clearly separated from the ring.
    rp = slide.get("right_panel")
    if rp:
        panel_x = Inches(8.0)
        panel_w = Inches(4.3)
        panel_y = Inches(1.7)
        mcore.add_text(s, panel_x, panel_y, panel_w, Inches(0.4), str(rp[0]),
                       font_size=Pt(18), font_color=mc.NAVY, bold=True)
        mcore.add_hline(s, panel_x, panel_y + Inches(0.45), panel_w, mc.NAVY, Pt(1))
        for i, pt in enumerate(rp[1]):
            mcore.add_text(s, panel_x + Inches(0.1), panel_y + Inches(0.7) + Inches(0.5) * i,
                           panel_w - Inches(0.2), Inches(0.45),
                           "• " + str(pt), font_size=Pt(13), font_color=mc.DARK_GRAY)
    eng._footer(s, slide.get("source", ""))


def _render_image_bg(s, image_path: str) -> None:
    """Draw a full-bleed background image + a white scrim for legibility.

    Used by content layouts that carry an ``image`` field. Drawn at the END of
    a slide body via send-to-back so shapes stay on top.
    """
    import os
    import mck_ppt.core as mcore
    import mck_ppt.constants as mc
    from pptx.util import Inches
    if not image_path or not os.path.isfile(image_path):
        return
    pic = s.shapes.add_picture(image_path, 0, 0, mc.SW, mc.SH)
    # Move picture to back so all shapes render on top of it.
    sp = pic._element
    sp.getparent().remove(sp)
    s.shapes._spTree.insert(2, sp)
    # White scrim (semi-opaque) for text legibility.
    from pptx.enum.shapes import MSO_SHAPE
    scrim = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, mc.SW, mc.SH)
    scrim.fill.solid(); scrim.fill.fore_color.rgb = mc.WHITE
    # Alpha not directly settable via python-pptx; use a light fill by default.
    mcore._clean_shape(scrim)


