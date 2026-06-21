"""
PPT Ingest模块

扫描目录中的PPTX文件，对每页提取：
- slide PNG image (通过LibreOffice)
- extracted text
- slide object metadata JSON
"""
import os
import json
import re
import subprocess
from pptx import Presentation
from pptx.util import Emu


def _sanitize_name(name):
    """文件名安全化"""
    name = re.sub(r'[^\w\u4e00-\u9fff-]', '_', name)
    return name[:50]


def _make_slide_id(deck_name, slide_no):
    """生成唯一slide_id"""
    return f"{_sanitize_name(deck_name)}_slide_{slide_no:03d}"


def _extract_shape_info(shape, slide_width, slide_height):
    """提取单个shape的信息"""
    info = {
        "object_id": shape.shape_id,
        "name": shape.name,
        "type": str(shape.shape_type).split("(")[0].strip() if shape.shape_type else "unknown",
        "text": "",
        "x_inches": round(shape.left / 914400, 3) if shape.left else 0,
        "y_inches": round(shape.top / 914400, 3) if shape.top else 0,
        "w_inches": round(shape.width / 914400, 3) if shape.width else 0,
        "h_inches": round(shape.height / 914400, 3) if shape.height else 0,
        "x_ratio": round(shape.left / slide_width, 4) if shape.left and slide_width else 0,
        "y_ratio": round(shape.top / slide_height, 4) if shape.top and slide_height else 0,
        "w_ratio": round(shape.width / slide_width, 4) if shape.width and slide_width else 0,
        "h_ratio": round(shape.height / slide_height, 4) if shape.height and slide_height else 0,
        "font_family": None,
        "font_size": None,
        "font_bold": None,
        "font_color": None,
        "fill_color": None,
        "paragraph_count": 0,
        "bullet_count": 0,
    }

    try:
        if hasattr(shape, 'text') and shape.text:
            info["text"] = shape.text[:500]
    except:
        pass

    try:
        if shape.has_text_frame:
            info["paragraph_count"] = len(shape.text_frame.paragraphs)
            for para in shape.text_frame.paragraphs:
                if para.level and para.level > 0:
                    info["bullet_count"] += 1
                for run in para.runs:
                    if run.font.name:
                        info["font_family"] = run.font.name
                    if run.font.size:
                        info["font_size"] = round(run.font.size / 12700, 1)  # to pt
                    info["font_bold"] = run.font.bold
                    try:
                        if run.font.color and run.font.color.rgb:
                            info["font_color"] = str(run.font.color.rgb)
                    except:
                        pass
                    break
                break
    except:
        pass

    try:
        if shape.fill and shape.fill.type is not None:
            try:
                info["fill_color"] = str(shape.fill.fore_color.rgb)
            except:
                pass
    except:
        pass

    return info


def _compute_features(objects, slide_width, slide_height):
    """计算页面特征"""
    features = {
        "text_shape_count": 0,
        "shape_count": 0,
        "picture_count": 0,
        "table_count": 0,
        "connector_count": 0,
        "group_count": 0,
        "image_area_ratio": 0.0,
        "total_text_chars": 0,
        "likely_image_slide": False,
    }

    total_image_area = 0.0
    slide_area = (slide_width / 914400) * (slide_height / 914400) if slide_width and slide_height else 1.0

    for obj in objects:
        obj_type = obj.get("type", "").upper()
        if "PICTURE" in obj_type:
            features["picture_count"] += 1
            total_image_area += obj["w_inches"] * obj["h_inches"]
        elif "TABLE" in obj_type:
            features["table_count"] += 1
        elif "LINE" in obj_type or "CONNECTOR" in obj_type:
            features["connector_count"] += 1
        elif "GROUP" in obj_type:
            features["group_count"] += 1
        else:
            features["shape_count"] += 1

        if obj.get("text"):
            features["text_shape_count"] += 1
            features["total_text_chars"] += len(obj["text"])

    features["image_area_ratio"] = round(total_image_area / slide_area, 3) if slide_area > 0 else 0
    features["likely_image_slide"] = (
        features["image_area_ratio"] > 0.7 and features["text_shape_count"] < 3
    )

    return features


def _render_slides(pptx_path, output_dir):
    """用LibreOffice将PPTX转PDF再转PNG"""
    try:
        # 转PDF (输出到临时目录，避免污染原目录)
        result = subprocess.run(
            ["soffice", "--headless", "--convert-to", "pdf", "--outdir", output_dir, pptx_path],
            capture_output=True, text=True, timeout=120,
        )
        if result.returncode != 0:
            print(f"  ⚠ LibreOffice conversion failed: {result.stderr[:200]}")
            return []

        pdf_path = os.path.join(output_dir, os.path.splitext(os.path.basename(pptx_path))[0] + ".pdf")
        if not os.path.exists(pdf_path):
            print(f"  ⚠ PDF not found after conversion")
            return []

        # 转PNG
        result = subprocess.run(
            ["pdftoppm", "-png", "-r", "150", pdf_path, os.path.join(output_dir, "page")],
            capture_output=True, text=True, timeout=120,
        )

        # 收集生成的PNG
        pngs = sorted([f for f in os.listdir(output_dir) if f.startswith("page") and f.endswith(".png")])

        # 清理PDF
        try:
            os.unlink(pdf_path)
        except:
            pass

        return pngs

    except Exception as e:
        print(f"  ⚠ Render failed: {e}")
        return []


def ingest_pptx_directory(ppt_dir, output_dir):
    """主入口：扫描目录，提取所有PPT的信息"""
    os.makedirs(output_dir, exist_ok=True)
    os.makedirs(os.path.join(output_dir, "slide_images"), exist_ok=True)
    os.makedirs(os.path.join(output_dir, "slide_objects"), exist_ok=True)
    os.makedirs(os.path.join(output_dir, "slide_text"), exist_ok=True)
    os.makedirs(os.path.join(output_dir, "slide_summaries"), exist_ok=True)
    os.makedirs(os.path.join(output_dir, "aggregate"), exist_ok=True)

    # 扫描PPTX文件
    pptx_files = []
    for f in os.listdir(ppt_dir):
        if f.endswith('.pptx') and not f.startswith('~'):
            pptx_files.append(os.path.join(ppt_dir, f))

    if not pptx_files:
        print(f"No .pptx files found in {ppt_dir}")
        return

    print(f"Found {len(pptx_files)} PPTX files")

    corpus_index = {
        "total_decks": len(pptx_files),
        "total_slides": 0,
        "slides": [],
        "errors": [],
    }

    for pptx_path in sorted(pptx_files):
        deck_name = os.path.splitext(os.path.basename(pptx_path))[0]
        print(f"\nProcessing: {deck_name}")

        try:
            prs = Presentation(pptx_path)
            slide_width = prs.slide_width
            slide_height = prs.slide_height

            # 渲染截图
            import tempfile
            with tempfile.TemporaryDirectory() as tmpdir:
                pngs = _render_slides(pptx_path, tmpdir)

                for slide_idx, slide in enumerate(prs.slides):
                    slide_no = slide_idx + 1
                    slide_id = _make_slide_id(deck_name, slide_no)

                    # 提取对象
                    objects = []
                    for shape in slide.shapes:
                        try:
                            obj = _extract_shape_info(shape, slide_width, slide_height)
                            objects.append(obj)
                        except Exception as e:
                            objects.append({"object_id": getattr(shape, 'shape_id', 0),
                                           "type": "error", "error": str(e)})

                    features = _compute_features(objects, slide_width, slide_height)

                    # 提取文字
                    all_text = []
                    for shape in slide.shapes:
                        try:
                            if hasattr(shape, 'text') and shape.text:
                                all_text.append(shape.text)
                        except:
                            pass

                    # 保存对象JSON
                    obj_data = {
                        "slide_id": slide_id,
                        "deck_name": deck_name,
                        "slide_no": slide_no,
                        "slide_size": {
                            "width": round(slide_width / 914400, 3),
                            "height": round(slide_height / 914400, 3),
                        },
                        "objects": objects,
                        "features": features,
                    }
                    obj_path = os.path.join(output_dir, "slide_objects", f"{slide_id}.json")
                    with open(obj_path, 'w', encoding='utf-8') as f:
                        json.dump(obj_data, f, ensure_ascii=False, indent=2)

                    # 保存文字
                    text_path = os.path.join(output_dir, "slide_text", f"{slide_id}.txt")
                    with open(text_path, 'w', encoding='utf-8') as f:
                        f.write("\n---\n".join(all_text))

                    # 复制截图
                    if slide_idx < len(pngs):
                        src_png = os.path.join(tmpdir, pngs[slide_idx])
                        dst_png = os.path.join(output_dir, "slide_images", f"{slide_id}.png")
                        import shutil
                        shutil.copy2(src_png, dst_png)

                    corpus_index["slides"].append({
                        "slide_id": slide_id,
                        "deck_name": deck_name,
                        "slide_no": slide_no,
                        "features": features,
                        "has_image": slide_idx < len(pngs),
                    })
                    corpus_index["total_slides"] += 1

            print(f"  ✓ {len(prs.slides)} slides extracted")

        except Exception as e:
            print(f"  ✗ Error: {e}")
            corpus_index["errors"].append({"deck": deck_name, "error": str(e)})

    # 保存索引
    index_path = os.path.join(output_dir, "corpus_index.json")
    with open(index_path, 'w', encoding='utf-8') as f:
        json.dump(corpus_index, f, ensure_ascii=False, indent=2)

    print(f"\n=== Ingest完成 ===")
    print(f"总计: {corpus_index['total_decks']} decks, {corpus_index['total_slides']} slides")
    print(f"索引: {index_path}")
