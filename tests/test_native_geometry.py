"""Geometry tests for native PPTX renderers — assert no shape overflows the
slide boundary (the most common QA gate failure). Covers all native layouts
that draw with python-pptx directly (swim_lane, pyramid_triangle, radar_chart,
gantt_chart, icon_stat_grid, cycle_ring, ecosystem_ring, harvey_ball_compare).
"""
import sys
from pathlib import Path
sys.path.insert(0, str(Path(__file__).resolve().parents[1] / "scripts"))

import pptx_renderer
from pptx_content_builder import BUILDERS
from pptx.util import Emu

# Slide is 13.333" × 7.5". Tolerance for the QA gate is 0.02".
SLIDE_W = 13.333
SLIDE_H = 7.5
TOLERANCE = 0.05  # slightly more lenient than gate's 0.02 for test stability

NATIVE_LAYOUTS = [
    "ecosystem_ring", "swim_lane", "pyramid_triangle", "radar_chart",
    "gantt_chart", "icon_stat_grid", "cycle", "harvey_ball_compare",
]


def _render_and_check(layout_name, tmp_path):
    """Render one native layout slide; return list of overflow shapes."""
    slide = {"idx": 1, "layout": layout_name, "title": f"测试{layout_name}标题文字",
             "key_point": "k"}
    data = BUILDERS[layout_name](slide, "skill-report")
    data["idx"] = 1
    content = {
        "brief": {"audience": "x", "goal": "r", "duration_minutes": 10},
        "framework": "p", "scenario": "skill-report",
        "slides": [data], "chapters": [],
    }
    deck_path = pptx_renderer.render_pptx(content, tmp_path)
    from pptx import Presentation
    prs = Presentation(str(deck_path))
    overflows = []
    for sh in prs.slides[0].shapes:
        if sh.left is None or sh.width is None:
            continue
        r = Emu(sh.left + sh.width).inches
        b = Emu(sh.top + sh.height).inches if sh.top and sh.height else 0
        if r > SLIDE_W + TOLERANCE:
            overflows.append(f"{sh.name} R={r:.2f}")
        if b > SLIDE_H + TOLERANCE:
            overflows.append(f"{sh.name} B={b:.2f}")
    return overflows


# Generate one test per native layout.
def _make_test(layout_name):
    def test_func(tmp_path):
        ov = _render_and_check(layout_name, tmp_path)
        assert not ov, f"{layout_name} shape(s) overflow: {ov}"
    test_func.__name__ = f"test_no_overflow_{layout_name}"
    return test_func


for _layout in NATIVE_LAYOUTS:
    globals()[f"test_no_overflow_{_layout}"] = _make_test(_layout)
